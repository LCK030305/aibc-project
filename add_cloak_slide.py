"""Insert a NEW slide right after the CrewAI 15-agent architecture slide
explaining CLOAK — Singapore's WOG PII sanitiser used at Stage 0.

Three sections on one slide:
  1. Public / Private key handshake — with the wax-seal analogy
  2. 7-entity redaction table (Singapore-tuned)
  3. Fail-CLOSED design principle (Topic 2.7)

Idempotent — if a slide with the CLOAK marker already exists it is
removed before the fresh one is inserted. Auto-detects the current
position of the CrewAI slide so this script keeps working even if
slide order changes.

Run:
    python add_cloak_slide.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu

PPT = Path(__file__).parent / (
    "AI Champions Bootcamp! Capstone project_1 of 2_Lim Chee Kuen (Danny)_"
    "1998459F_May 2026.pptx"
)

SLIDE_TITLE = "CLOAK — Stage 0 PII Protection"
SLIDE_MARKER = "How the key handshake works"  # unique to this slide
CREWAI_MARKER = "Hierarchical Crew"  # identifies the CrewAI arch slide

# --- Palette ------------------------------------------------------------
C_TITLE   = RGBColor(0x1F, 0x2A, 0x44)  # navy
C_TEXT    = RGBColor(0x11, 0x18, 0x27)  # near-black
C_MUTED   = RGBColor(0x6B, 0x72, 0x80)  # grey (caption / footer)

# Public key — safe / visible (blue)
C_PUB_FILL  = RGBColor(0xDB, 0xEA, 0xFE)
C_PUB_LINE  = RGBColor(0x3B, 0x82, 0xF6)
C_PUB_TEXT  = RGBColor(0x1E, 0x3A, 0x8A)

# Private key — secret / danger-if-leaked (red)
C_PRIV_FILL = RGBColor(0xFE, 0xE2, 0xE2)
C_PRIV_LINE = RGBColor(0xEF, 0x44, 0x44)
C_PRIV_TEXT = RGBColor(0x7F, 0x1D, 0x1D)

# Fail-CLOSED callout (amber — matches other slides)
C_CALL_FILL = RGBColor(0xFE, 0xF3, 0xC7)
C_CALL_LINE = RGBColor(0xF5, 0x9E, 0x0B)
C_CALL_TEXT = RGBColor(0x78, 0x35, 0x0F)

# Analogy strip (neutral)
C_STRIP_FILL = RGBColor(0xF9, 0xFA, 0xFB)
C_STRIP_LINE = RGBColor(0xE5, 0xE7, 0xEB)

# Table row alternation
C_TABLE_HEAD_FILL = RGBColor(0x1F, 0x2A, 0x44)
C_TABLE_HEAD_TEXT = RGBColor(0xFF, 0xFF, 0xFF)
C_TABLE_ODD_FILL  = RGBColor(0xFF, 0xFF, 0xFF)
C_TABLE_EVEN_FILL = RGBColor(0xF9, 0xFA, 0xFB)


def _set_text(shape, text, *, size=11, bold=False, color=C_TEXT,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
              italic=False, mono=False):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    lines = text.split("\n")
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = lines[0]
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    if mono:
        r.font.name = "Consolas"
    for line in lines[1:]:
        ep = tf.add_paragraph(); ep.alignment = align
        r2 = ep.add_run(); r2.text = line
        r2.font.size = Pt(size); r2.font.bold = bold; r2.font.italic = italic
        r2.font.color.rgb = color
        if mono:
            r2.font.name = "Consolas"


def _remove_slide_at_index(prs, idx):
    sldIdLst = prs.slides._sldIdLst
    sldId = list(sldIdLst)[idx]
    rId = sldId.attrib[qn("r:id")]
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)


def _move_slide(prs, old_index, new_index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(new_index, slides[old_index])


def _slide_has_marker(slide, marker: str) -> bool:
    for sh in slide.shapes:
        if sh.has_text_frame and marker in sh.text_frame.text:
            return True
    return False


def _find_slide_index_with_marker(prs, marker: str) -> int:
    for i, s in enumerate(prs.slides):
        if _slide_has_marker(s, marker):
            return i
    return -1


def build_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    SW = prs.slide_width / 914400  # 13.33"

    # --- Title ---
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.30), Inches(SW - 1.0), Inches(0.55),
    )
    _set_text(
        title_box, SLIDE_TITLE,
        size=26, bold=True, color=C_TITLE,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
    )

    # --- Subtitle ---
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.88), Inches(SW - 1.0), Inches(0.28),
    )
    _set_text(
        subtitle_box,
        "Singapore's Whole-of-Government privacy toolkit · Topic 5.5.2",
        size=12, italic=True, color=C_MUTED,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
    )

    # ================================================================
    # SECTION 1 — Public / Private key handshake
    # ================================================================
    SEC1_Y = 1.25
    # Section header
    hdr1 = slide.shapes.add_textbox(
        Inches(0.5), Inches(SEC1_Y), Inches(SW - 1.0), Inches(0.32),
    )
    _set_text(
        hdr1, "🔐  How the key handshake works",
        size=14, bold=True, color=C_TITLE,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
    )

    # Two key boxes side-by-side
    BOX_Y = SEC1_Y + 0.40
    BOX_H = 1.10
    MARGIN = 0.5
    GAP = 0.20
    BOX_W = (SW - 2 * MARGIN - GAP) / 2

    # LEFT — PUBLIC KEY
    pub = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(MARGIN), Inches(BOX_Y), Inches(BOX_W), Inches(BOX_H),
    )
    pub.fill.solid(); pub.fill.fore_color.rgb = C_PUB_FILL
    pub.line.color.rgb = C_PUB_LINE; pub.line.width = Pt(1.5)
    pub.shadow.inherit = False
    _set_text(
        pub,
        "🌐  PUBLIC KEY\n"
        "Like your username — sent in every request header.\n"
        "Identifies WHO the request is from. Safe if intercepted.",
        size=12, bold=False, color=C_PUB_TEXT,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    # Make the first line bold+bigger
    tf = pub.text_frame
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.size = Pt(15)

    # RIGHT — PRIVATE KEY
    priv = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(MARGIN + BOX_W + GAP), Inches(BOX_Y),
        Inches(BOX_W), Inches(BOX_H),
    )
    priv.fill.solid(); priv.fill.fore_color.rgb = C_PRIV_FILL
    priv.line.color.rgb = C_PRIV_LINE; priv.line.width = Pt(1.5)
    priv.shadow.inherit = False
    _set_text(
        priv,
        "🔒  PRIVATE KEY\n"
        "Like your password — NEVER sent over the network.\n"
        "Stays local; used to sign each request with a unique fingerprint.",
        size=12, bold=False, color=C_PRIV_TEXT,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    tf = priv.text_frame
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.size = Pt(15)

    # Analogy strip below
    STRIP_Y = BOX_Y + BOX_H + 0.10
    STRIP_H = 0.48
    strip = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(MARGIN), Inches(STRIP_Y),
        Inches(SW - 2 * MARGIN), Inches(STRIP_H),
    )
    strip.fill.solid(); strip.fill.fore_color.rgb = C_STRIP_FILL
    strip.line.color.rgb = C_STRIP_LINE; strip.line.width = Pt(0.75)
    strip.shadow.inherit = False
    _set_text(
        strip,
        "🕯️  Think of a wax seal: only your ring makes it, it cracks "
        "if the letter changes, and it expires at midnight.",
        size=12, bold=False, italic=True, color=C_TEXT,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )

    # ================================================================
    # SECTION 2 — 7-entity redaction table
    # ================================================================
    SEC2_Y = STRIP_Y + STRIP_H + 0.20
    hdr2 = slide.shapes.add_textbox(
        Inches(0.5), Inches(SEC2_Y), Inches(SW - 1.0), Inches(0.30),
    )
    _set_text(
        hdr2, "🛡️  What CLOAK redacts (7 Singapore-tuned entities)",
        size=14, bold=True, color=C_TITLE,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
    )

    TBL_Y = SEC2_Y + 0.35
    TBL_H = 1.65
    TBL_W = SW - 2 * MARGIN

    entities = [
        ("Entity type",       "Original example",       "Replaced with"),
        ("PERSON",             "Mdm Kim Harin",          "<PERSON>"),
        ("SG_NRIC_FIN",        "S8273756Y",              "<SG_NRIC_FIN>"),
        ("PHONE_NUMBER",       "8123 4567",              "<PHONE_NUMBER>"),
        ("EMAIL_ADDRESS",      "kim@gmail.com",          "<EMAIL_ADDRESS>"),
        ("SG_BANK_ACCOUNT_NO", "DBS 011-12345-6",        "<SG_BANK_ACCOUNT_NO>"),
        ("SG_ADDRESS",         "Blk 25 Sin Ming Rd",     "<SG_ADDRESS>"),
        ("DATE_TIME",          "born 3 April 1978",      "<DATE_TIME>"),
    ]

    table_shape = slide.shapes.add_table(
        rows=len(entities), cols=3,
        left=Inches(MARGIN), top=Inches(TBL_Y),
        width=Inches(TBL_W), height=Inches(TBL_H),
    )
    tbl = table_shape.table
    # Column widths (roughly proportional)
    tbl.columns[0].width = Inches(TBL_W * 0.28)
    tbl.columns[1].width = Inches(TBL_W * 0.32)
    tbl.columns[2].width = Inches(TBL_W * 0.40)

    for r_idx, row in enumerate(entities):
        for c_idx, val in enumerate(row):
            cell = tbl.cell(r_idx, c_idx)
            cell.margin_left = cell.margin_right = Inches(0.10)
            cell.margin_top = cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            tf_cell = cell.text_frame
            tf_cell.clear()
            tf_cell.word_wrap = True
            p_cell = tf_cell.paragraphs[0]
            p_cell.alignment = PP_ALIGN.LEFT
            run = p_cell.add_run(); run.text = val

            if r_idx == 0:  # header
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_TABLE_HEAD_FILL
                run.font.size = Pt(11)
                run.font.bold = True
                run.font.color.rgb = C_TABLE_HEAD_TEXT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = (
                    C_TABLE_EVEN_FILL if r_idx % 2 == 0 else C_TABLE_ODD_FILL
                )
                run.font.size = Pt(10)
                run.font.color.rgb = C_TEXT
                # Entity type & redacted token in monospace
                if c_idx == 0 or c_idx == 2:
                    run.font.name = "Consolas"
                    if c_idx == 2:
                        run.font.color.rgb = C_PRIV_TEXT

    # ================================================================
    # SECTION 3 — Fail-CLOSED callout
    # ================================================================
    SEC3_Y = TBL_Y + TBL_H + 0.15
    SEC3_H = 0.85
    fc = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(MARGIN), Inches(SEC3_Y),
        Inches(SW - 2 * MARGIN), Inches(SEC3_H),
    )
    fc.fill.solid(); fc.fill.fore_color.rgb = C_CALL_FILL
    fc.line.color.rgb = C_CALL_LINE; fc.line.width = Pt(1.5)
    fc.shadow.inherit = False
    _set_text(
        fc,
        "⚠️  Fail-CLOSED (Topic 2.7): if CLOAK is unreachable or keys "
        "are missing, the request FAILS — raw PII is never passed "
        "through to the LLM.",
        size=13, bold=True, color=C_CALL_TEXT,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )

    # --- Footer ---
    foot = slide.shapes.add_textbox(
        Inches(0.5), Inches(7.10), Inches(SW - 1.0), Inches(0.30),
    )
    _set_text(
        foot,
        "1 CLOAK call per query · ~300 ms latency · $0 cost · "
        "fixed at Stage 0 of the pipeline",
        size=9, italic=True, color=C_MUTED,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )


def main():
    prs = Presentation(PPT)
    total_before = len(prs.slides)
    print(f"Opened PPT — {total_before} slides before.")

    # Idempotency: remove any existing CLOAK slide first
    for idx, sl in enumerate(prs.slides):
        if _slide_has_marker(sl, SLIDE_MARKER):
            _remove_slide_at_index(prs, idx)
            print(f"Removed existing CLOAK slide (was slide {idx + 1}).")
            break

    # Find CrewAI slide dynamically
    crewai_idx = _find_slide_index_with_marker(prs, CREWAI_MARKER)
    if crewai_idx < 0:
        raise RuntimeError(
            f"Could not find CrewAI slide (marker: {CREWAI_MARKER!r}). "
            "Update CREWAI_MARKER in this script."
        )
    target_idx = crewai_idx + 1
    print(f"CrewAI slide at position {crewai_idx + 1}; "
          f"CLOAK will go at position {target_idx + 1}.")

    # Build and move into place
    build_slide(prs)
    new_index = len(prs.slides) - 1
    _move_slide(prs, new_index, target_idx)
    print(f"Moved new slide from position {new_index + 1} "
          f"to position {target_idx + 1}.")

    prs.save(PPT)
    print(f"Saved — {len(prs.slides)} slides total.")


if __name__ == "__main__":
    main()
