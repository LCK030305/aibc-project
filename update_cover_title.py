"""Update the cover slide (slide 1) title of Danny's Capstone PPT.

Old title  : "Improve productivity of Social Assistance Officers in
              identifying relevant welfare resources of needy citizens"
New title  : "Matching Welfare Resources to the Families They Serve"
Old sub    : "Capstone Project (use case 1 of 2)"
New sub    : "AI-Powered Decision Support for Social Assistance Officers"

Preserves the existing formatting (font, size, color) of the shapes
by replacing only the text content of the first run in each shape's
first paragraph, and removing any leftover runs / paragraphs.

Idempotent — running twice produces the same result.
"""

from pathlib import Path
from pptx import Presentation

PPT = Path(__file__).parent / (
    "AI Champions Bootcamp! Capstone project_1 of 2_Lim Chee Kuen (Danny)_"
    "1998459F_May 2026.pptx"
)

NEW_TITLE = "Matching Welfare Resources to the Families They Serve"
NEW_SUBTITLE = "Multi-Agent Decision Support for Social Assistance Officers"


def _replace_text_preserve_format(shape, new_text: str) -> None:
    """Replace shape's text but keep the first-run formatting intact."""
    tf = shape.text_frame
    if not tf.paragraphs or not tf.paragraphs[0].runs:
        # Fallback if the paragraph has no runs — clear + add plain
        tf.clear()
        tf.paragraphs[0].add_run().text = new_text
        return

    first_para = tf.paragraphs[0]
    first_run = first_para.runs[0]
    first_run.text = new_text

    # Remove any additional runs in the first paragraph
    for run in list(first_para.runs)[1:]:
        run._r.getparent().remove(run._r)

    # Remove any additional paragraphs beyond the first
    for para in list(tf.paragraphs)[1:]:
        para._p.getparent().remove(para._p)


def main() -> None:
    prs = Presentation(PPT)
    slide = prs.slides[0]

    # From inspection: shape[1] is the title, shape[2] is the subtitle.
    title_shape = slide.shapes[1]
    subtitle_shape = slide.shapes[2]

    print(f"Old title    : {title_shape.text_frame.text!r}")
    print(f"Old subtitle : {subtitle_shape.text_frame.text!r}")

    _replace_text_preserve_format(title_shape, NEW_TITLE)
    _replace_text_preserve_format(subtitle_shape, NEW_SUBTITLE)

    print(f"New title    : {title_shape.text_frame.text!r}")
    print(f"New subtitle : {subtitle_shape.text_frame.text!r}")

    prs.save(PPT)
    print("Saved.")


if __name__ == "__main__":
    main()
