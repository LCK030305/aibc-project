"""Insert a NEW slide at position 3 (right after 'About MSF') of Danny's
Capstone PPT, showing a 3-line problem statement distilled from the
full 150-word paragraph.

Structure:
  - Title:      "Problem Statement"
  - Subtitle:   "The gap this capstone closes"
  - 3 lines:
      Line 1 (situation): 750+ schemes across 12 categories
      Line 2 (impact):    hours per case → less time for engagement
      Line 3 (HMW):       amber-highlighted callout — the driving question
  - Footer:     source attribution

Idempotent — if a slide with the marker phrase already exists at
position 3, it's removed before the fresh one is inserted.

Run:
    python add_problem_slide.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

PPT = Path(__file__).parent / (
    "AI Champions Bootcamp! Capstone project_1 of 2_Lim Chee Kuen (Danny)_"
    "1998459F_May 2026.pptx"
)

SLIDE_TITLE = "Problem Statement"
SLIDE_MARKER = "How might we surface the right resources"  # unique phrase

LINE_1 = (
    "📋  SAOs at MSF sift through 750+ welfare schemes across 12 "
    "categories to match each client's unique needs."
)
LINE_2 = (
    "⏱️  Manual cross-referencing takes hours per case — shrinking "
    "time for real-time conversations and meaningful client engagement."
)
LINE_3 = (
    "🎯  How might we surface the right resources instantly, so "
    "SAOs can focus on the families they serve?"
)

C_TITLE = RGBColor(0x1F, 0x2A, 0x44)
C_TEXT = RGBColor(0x11, 0x18, 0x27)
C_MUTED = RGBColor(0x6B, 0x72, 0x80)
C_LINE_FILL = RGBColor(0xF9, 0xFA, 0xFB)
C_LINE_LINE = RGBColor(0xE5, 0xE7, 0xEB)
C_CALLOUT_FILL = RGBColor(0xFE, 0xF3, 0xC7)
C_CALLOUT_LINE = RGBColor(0xF5, 0x9E, 0x0B)
C_CALLOUT_TEXT = RGBColor(0x78, 0x35, 0x0F)


def _set_text(shape, text, *, size=11, bold=False, color=C_TEXT,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
              italic=False):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.15)
    tf.margin_top = tf.margin_bottom = Inches(0.08)
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color


def _remove_slide_at_index(prs, idx):
    sldIdLst = prs.slides._sldIdLst
    sldId = list(sldIdLst)[idx]
    rId = sldId.attrib[qn("r:id")]
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)


def _move_slide(prs, old_index, new_index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(new_index, slides[old_index])


def _slide_has_marker(slide, marker: str) -> bool:
    for sh in slide.shapes:
        if sh.has_text_frame and marker in sh.text_frame.text:
            return True
    return False


def build_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    SW = prs.slide_width / 914400  # 13.33"

    # --- Title ---
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.30), Inches(SW - 1.0), Inches(0.55),
    )
    _set_text(
        title_box, SLIDE_TITLE,
        size=28, bold=True, color=C_TITLE,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
    )

    # --- Subtitle ---
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.92), Inches(SW - 1.0), Inches(0.30),
    )
    _set_text(
        subtitle_box, "The gap this capstone closes",
        size=13, italic=True, color=C_MUTED,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
    )

    # --- Body: Line 1 (situation) ---
    L1_Y = 1.90
    L1_H = 1.10
    line1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.7), Inches(L1_Y), Inches(SW - 1.4), Inches(L1_H),
    )
    line1.fill.solid(); line1.fill.fore_color.rgb = C_LINE_FILL
    line1.line.color.rgb = C_LINE_LINE; line1.line.width = Pt(0.75)
    line1.shadow.inherit = False
    _set_text(
        line1, LINE_1,
        size=20, bold=False, color=C_TEXT,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
    )

    # --- Body: Line 2 (impact) ---
    L2_Y = L1_Y + L1_H + 0.25
    L2_H = 1.10
    line2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.7), Inches(L2_Y), Inches(SW - 1.4), Inches(L2_H),
    )
    line2.fill.solid(); line2.fill.fore_color.rgb = C_LINE_FILL
    line2.line.color.rgb = C_LINE_LINE; line2.line.width = Pt(0.75)
    line2.shadow.inherit = False
    _set_text(
        line2, LINE_2,
        size=20, bold=False, color=C_TEXT,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
    )

    # --- Body: Line 3 (HMW callout — amber highlight) ---
    L3_Y = L2_Y + L2_H + 0.25
    L3_H = 1.20
    line3 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.7), Inches(L3_Y), Inches(SW - 1.4), Inches(L3_H),
    )
    line3.fill.solid(); line3.fill.fore_color.rgb = C_CALLOUT_FILL
    line3.line.color.rgb = C_CALLOUT_LINE; line3.line.width = Pt(1.5)
    line3.shadow.inherit = False
    _set_text(
        line3, LINE_3,
        size=22, bold=True, color=C_CALLOUT_TEXT,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
    )

    # --- Footer ---
    foot = slide.shapes.add_textbox(
        Inches(0.5), Inches(7.10), Inches(SW - 1.0), Inches(0.30),
    )
    _set_text(
        foot,
        "Distilled from the full 150-word problem-statement paragraph "
        "(see samples/problem statement.jpeg for the AI-assisted original)",
        size=9, italic=True, color=C_MUTED,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )


def main():
    prs = Presentation(PPT)
    total_before = len(prs.slides)
    print(f"Opened PPT — {total_before} slides before.")

    # Idempotency: if any existing slide has our marker, remove it
    for idx, sl in enumerate(prs.slides):
        if _slide_has_marker(sl, SLIDE_MARKER):
            _remove_slide_at_index(prs, idx)
            print(f"Removed existing Problem Statement slide "
                  f"(was slide {idx + 1}).")
            break

    build_slide(prs)
    new_index = len(prs.slides) - 1
    _move_slide(prs, new_index, 2)  # position 3 → index 2
    print(f"Moved new slide from position {new_index + 1} to position 3.")

    prs.save(PPT)
    print(f"Saved — {len(prs.slides)} slides total.")


if __name__ == "__main__":
    main()
