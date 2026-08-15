"""Insert a NEW slide at position 2 (right after the cover) of Danny's
Capstone PPT, showing MSF's 'What We Do' screenshot as context for the
capstone's use-case scope.

Layout:
  - Title:      "About MSF — Where Our Use Case Fits"
  - Caption:    "Source: Ministry of Social and Family Development (msf.gov.sg)"
  - Image:      msf_what_we_do.png (centred, 12" wide)
  - Callout:    highlights that this capstone addresses Financial
                Assistance & Social Support
  - Footer:     "Context: where the SAO Programme Matcher plugs into MSF's
                welfare service portfolio"

Idempotent — if a slide with the marker title already exists at position 2,
it is removed before the fresh one is inserted (safe re-run).

Run:
    python add_context_slide.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

PPT = Path(__file__).parent / (
    "AI Champions Bootcamp! Capstone project_1 of 2_Lim Chee Kuen (Danny)_"
    "1998459F_May 2026.pptx"
)
IMAGE = Path(__file__).parent / "samples" / "msf_what_we_do.png"

SLIDE_TITLE = "About MSF — Where Our Use Case Fits"
SLIDE_MARKER = "About MSF"  # used to detect prior insertion

C_TITLE = RGBColor(0x1F, 0x2A, 0x44)      # navy (matches cover)
C_TEXT = RGBColor(0x11, 0x18, 0x27)       # near-black
C_MUTED = RGBColor(0x6B, 0x72, 0x80)      # grey (caption / footer)
C_CALLOUT_FILL = RGBColor(0xFE, 0xF3, 0xC7)  # amber (matches highlighted tab)
C_CALLOUT_LINE = RGBColor(0xF5, 0x9E, 0x0B)
C_CALLOUT_TEXT = RGBColor(0x78, 0x35, 0x0F)  # dark amber
C_IMG_BORDER = RGBColor(0xE5, 0xE7, 0xEB)


def _set_text(shape, text, *, size=11, bold=False, color=C_TEXT,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
              italic=False):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]; p.alignment = align
    lines = text.split("\n")
    run = p.add_run(); run.text = lines[0]
    run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color
    for line in lines[1:]:
        ep = tf.add_paragraph(); ep.alignment = align
        r2 = ep.add_run(); r2.text = line
        r2.font.size = Pt(size); r2.font.bold = bold; r2.font.italic = italic
        r2.font.color.rgb = color


def _remove_slide_at_index(prs, idx):
    """Remove slide at `idx` from the presentation (0-indexed)."""
    sldIdLst = prs.slides._sldIdLst
    sldId = list(sldIdLst)[idx]
    rId = sldId.attrib[qn("r:id")]
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)


def _move_slide(prs, old_index, new_index):
    """Reorder slides: move slide from old_index to new_index."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(new_index, slides[old_index])


def _slide_has_marker(slide, marker: str) -> bool:
    for sh in slide.shapes:
        if sh.has_text_frame and marker in sh.text_frame.text:
            return True
    return False


def build_slide(prs):
    blank_layout = prs.slide_layouts[6]  # "Blank"
    slide = prs.slides.add_slide(blank_layout)
    SW = prs.slide_width / 914400  # 13.33"
    CX = SW / 2

    # --- Title ---
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.30), Inches(SW - 1.0), Inches(0.55),
    )
    _set_text(
        title_box, SLIDE_TITLE,
        size=26, bold=True, color=C_TITLE,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
    )

    # --- Caption (source attribution) ---
    caption_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.90), Inches(SW - 1.0), Inches(0.30),
    )
    _set_text(
        caption_box,
        "Source: Ministry of Social and Family Development (msf.gov.sg)",
        size=11, italic=True, color=C_MUTED,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
    )

    # --- Image (centred, 12" wide, aspect-ratio preserved) ---
    IMG_W = 12.0
    img_x = CX - IMG_W / 2
    img_y = 1.55
    pic = slide.shapes.add_picture(
        str(IMAGE), Inches(img_x), Inches(img_y), width=Inches(IMG_W),
    )
    # Add a subtle border to the image
    pic.line.color.rgb = C_IMG_BORDER
    pic.line.width = Pt(1.0)

    # Compute image bottom for subsequent element placement
    img_h_in = pic.height / 914400
    img_bottom = img_y + img_h_in

    # --- Callout (amber highlight linking image to our capstone scope) ---
    CALL_H = 0.72
    call_y = max(img_bottom + 0.35, 5.30)
    call_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.0), Inches(call_y), Inches(SW - 2.0), Inches(CALL_H),
    )
    call_box.fill.solid()
    call_box.fill.fore_color.rgb = C_CALLOUT_FILL
    call_box.line.color.rgb = C_CALLOUT_LINE
    call_box.line.width = Pt(1.25)
    call_box.shadow.inherit = False
    _set_text(
        call_box,
        "🎯  This capstone addresses Financial Assistance & Social Support "
        "— the highlighted tab within Help Those in Need",
        size=14, bold=True, color=C_CALLOUT_TEXT,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )

    # --- Footer ---
    foot = slide.shapes.add_textbox(
        Inches(0.5), Inches(7.10), Inches(SW - 1.0), Inches(0.30),
    )
    _set_text(
        foot,
        "Context: where the SAO Programme Matcher plugs into MSF's "
        "welfare service portfolio",
        size=9, italic=True, color=C_MUTED,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )


def main():
    prs = Presentation(PPT)
    total_before = len(prs.slides)
    print(f"Opened PPT — {total_before} slides before.")

    # Idempotency: if current slide 2 already has our marker, remove it.
    if total_before >= 2 and _slide_has_marker(prs.slides[1], SLIDE_MARKER):
        _remove_slide_at_index(prs, 1)
        print("Removed existing 'About MSF' slide at position 2.")

    # Append new slide (goes to end)
    build_slide(prs)

    # Move it to position 2 (index 1 — right after cover at index 0)
    new_index = len(prs.slides) - 1
    _move_slide(prs, new_index, 1)
    print(f"Moved new slide from position {new_index + 1} to position 2.")

    prs.save(PPT)
    print(f"Saved — {len(prs.slides)} slides total.")


if __name__ == "__main__":
    main()
