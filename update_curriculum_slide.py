"""Update PPT slide 15 (Bootcamp curriculum coverage) with an expanded
table reflecting the current state of the solution.

Adds rows for:
  - Topic 2.4/2.5 Reasoning (Chain-of-Thought · Least-to-Most)
  - Topic 2.6 Human-in-the-Loop (3 HITL gates: case breakdown · top-5 ·
    case documentation)
  - Topic 2.6 Decision Chain (Safety → Router)
  - Topic 2.7 Exception Handling (fail-CLOSED CLOAK + safety)
  - Topic 4.4 Post-Retrieval Faithfulness audit
  - Topic 5.5 CrewAI multi-agent (15 agents · Deep Mode)
  - Topic 5.5.2 CLOAK PII Free-Text Anonymisation

Idempotent — replaces the whole table on slide 15 with the new one.

Run:
    python update_curriculum_slide.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

PPT = Path(__file__).parent / (
    "AI Champions Bootcamp! Capstone project_1 of 2_Lim Chee Kuen (Danny)_"
    "1998459F_May 2026.pptx"
)

SLIDE_INDEX = 14  # slide 15 (0-indexed) — used as fallback only
SLIDE_TITLE = "Bootcamp curriculum coverage"
SLIDE_SUBTITLE = (
    "Every topic block the bootcamp teaches has a corresponding "
    "artefact in the repo."
)
SLIDE_FOOTER = "Capstone deadline: 14 August 2026 · Politemall"

C_HEAD_FILL = RGBColor(0x1F, 0x2A, 0x44)
C_HEAD_TEXT = RGBColor(0xFF, 0xFF, 0xFF)
C_ROW_ODD   = RGBColor(0xFF, 0xFF, 0xFF)
C_ROW_EVEN  = RGBColor(0xF9, 0xFA, 0xFB)
C_TEXT      = RGBColor(0x11, 0x18, 0x27)
C_TOPIC     = RGBColor(0x1E, 0x3A, 0x8A)  # blue for the topic column

ROWS = [
    ("Week / Topic", "What the bootcamp teaches", "Where it lives in the Claude vibe coding repo"),
    ("Week 1", "LLM foundations · Prompt Engineering · CO-STAR · Tokens · Hallucinations",
     "llm.py · prompts.py · faithfulness_check.py"),
    ("Week 2 (2.4 / 2.5)", "Reasoning · Chain-of-Thought · Least-to-Most decomposition",
     "decomposer.py · recommender.py (reasoning_steps in JSON schema)"),
    ("Week 2 (2.6)", "Decision Chain: safety check · query router",
     "safety.py · router.py · recommender.py (Stage 1-2)"),
    ("Week 2 (2.6) — HITL", "Human-in-the-Loop as part of the workflow (3 gates)",
     "app.py: HITL #1 case breakdown · #2 top-5 (Agent #14) · #3 case docs (Agent #15)"),
    ("Week 2 (2.7)", "Exception handling · Fail-CLOSED design",
     "pii_filter.py · safety.py · recommender.py (rate-limit / auth guards)"),
    ("Week 3 (3.x)", "Embeddings · Vector search · RAG · Beyond keywords",
     "embed.py · retriever.py · bm25_retriever.py (hybrid) · hyde.py"),
    ("Week 4 (4.x)", "Advanced RAG · Pre / Post-retrieval · Faithfulness audit · Eval",
     "chunker.py · hyde.py · faithfulness_check.py · ragas_eval.py"),
    ("Week 5 (5.x)", "Agents · Multi-agent orchestration · Secure credentials",
     "crew_runner.py (CrewAI · 15 agents) · llm.py get_secret() · .env / st.secrets"),
    ("Week 5 (5.5)", "CrewAI multi-agent: Triaging · 12 specialists · Aggregator · Case Doc Officer",
     "crew_runner.py · agents_config.py · recommender.py _recommend_deep_mode()"),
    ("Week 5 (5.5.2)", "CLOAK Free-Text Anonymisation (WOG PII toolkit)",
     "pii_filter.py (Stage 0 · fail-CLOSED · CLOAK-AUTH HMAC signing)"),
    ("Week 6 (6.x)", "Streamlit basics · pip + venv · Debugging",
     "requirements.txt · .venv/ · app.py"),
    ("Week 8 (8.x)", "Streamlit Deep Dive · Password gate · Git · Deploy · Vibe coding",
     "app.py (@st.cache_resource, session_state) · Streamlit Cloud deployment"),
]


def _style_cell(cell, text: str, *, bold: bool = False,
                color: RGBColor = C_TEXT, size: int = 9,
                fill: RGBColor = None):
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    cell.margin_left = cell.margin_right = Inches(0.06)
    cell.margin_top = cell.margin_bottom = Inches(0.02)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _remove_shape(shape):
    shape._element.getparent().remove(shape._element)


def _find_slide_by_title(prs, title_marker):
    """Return the (index, slide) for the first slide whose text contains
    `title_marker`, else (None, None)."""
    for i, s in enumerate(prs.slides):
        for sh in s.shapes:
            if sh.has_text_frame and title_marker in sh.text_frame.text:
                return i, s
    return None, None


def _remove_slide_at_index(prs, idx):
    sldIdLst = prs.slides._sldIdLst
    sldId = list(sldIdLst)[idx]
    rId = sldId.attrib[qn("r:id")]
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)


def _build_curriculum_slide(prs):
    """Build the curriculum coverage slide from scratch, table + all."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    SW = prs.slide_width / 914400

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.35), Inches(SW - 1.0), Inches(0.60),
    )
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    run = p.add_run(); run.text = SLIDE_TITLE
    run.font.size = Pt(28); run.font.bold = True
    run.font.color.rgb = C_HEAD_FILL

    # Subtitle
    sub_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.00), Inches(SW - 1.0), Inches(0.35),
    )
    tf = sub_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    run = p.add_run(); run.text = SLIDE_SUBTITLE
    run.font.size = Pt(12); run.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)

    # Table
    n_rows = len(ROWS)
    n_cols = 3
    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.5), Inches(1.55), Inches(12.33), Inches(5.35),
    )
    tbl = table_shape.table
    tbl.columns[0].width = Inches(2.30)
    tbl.columns[1].width = Inches(4.60)
    tbl.columns[2].width = Inches(5.43)
    for r_idx, row in enumerate(ROWS):
        for c_idx, val in enumerate(row):
            cell = tbl.cell(r_idx, c_idx)
            if r_idx == 0:
                _style_cell(cell, val, bold=True, color=C_HEAD_TEXT,
                            size=10, fill=C_HEAD_FILL)
            else:
                fill = C_ROW_EVEN if r_idx % 2 == 0 else C_ROW_ODD
                color = C_TOPIC if c_idx == 0 else C_TEXT
                _style_cell(cell, val, bold=(c_idx == 0), color=color,
                            size=8, fill=fill)

    # Footer
    foot_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(7.10), Inches(SW - 1.0), Inches(0.30),
    )
    tf = foot_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    run = p.add_run(); run.text = SLIDE_FOOTER
    run.font.size = Pt(10); run.font.italic = True
    run.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)


def main():
    prs = Presentation(PPT)
    before = len(prs.slides)
    print(f"Opened PPT — {before} slides before.")

    # Find any existing curriculum slide (by title marker) and remove it
    orig_idx, _ = _find_slide_by_title(prs, SLIDE_TITLE)
    if orig_idx is not None:
        _remove_slide_at_index(prs, orig_idx)
        print(f"Removed existing curriculum slide (was slide {orig_idx + 1}).")
    else:
        orig_idx = len(prs.slides)  # append if not found
        print(f"No existing curriculum slide — will add as slide {orig_idx + 1}.")

    # Build fresh
    _build_curriculum_slide(prs)

    # Move to original position (or leave at end)
    new_index = len(prs.slides) - 1
    if new_index != orig_idx:
        sldIdLst = prs.slides._sldIdLst
        slides = list(sldIdLst)
        sldIdLst.remove(slides[new_index])
        sldIdLst.insert(orig_idx, slides[new_index])
        print(f"Moved from position {new_index + 1} to {orig_idx + 1}.")

    prs.save(PPT)
    print(f"Saved — {len(prs.slides)} slides after.")


if __name__ == "__main__":
    main()
