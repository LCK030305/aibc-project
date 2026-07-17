"""Add (or replace) a 'Deep Analysis Mode — CrewAI Hierarchical Crew'
architecture slide on Danny's Capstone PPT.

Renders the conceptual flow (SAO note → CLOAK → Safety → Crew container
→ Audit → HITL) with all 12 specialist agents shown in a 6×2 grid.

Idempotent — if the last slide is already this architecture slide
(detected by a marker in the title), it's removed before the fresh one
is added. Safe to re-run after tweaks.

Run:
    python add_crewai_slide.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from lxml import etree

PPT_FILE = Path(__file__).parent / (
    "AI Champions Bootcamp! Capstone project_1 of 2_Lim Chee Kuen (Danny)_"
    "1998459F_May 2026.pptx"
)

SLIDE_TITLE = (
    "15 Agents — Deep Analysis Mode · CrewAI Hierarchical Crew  (Topic 5.5)"
)
SLIDE_MARKER = "CrewAI Hierarchical Crew"  # used to detect prior insertion

# ─── Color palette ─────────────────────────────────────────────────────────
C_TITLE = RGBColor(0x1F, 0x2A, 0x44)
C_TEXT = RGBColor(0x11, 0x18, 0x27)
C_MUTED = RGBColor(0x6B, 0x72, 0x80)
C_UNCHANGED_FILL = RGBColor(0xEF, 0xF6, 0xFF)
C_UNCHANGED_LINE = RGBColor(0x3B, 0x82, 0xF6)
C_CREW_FILL = RGBColor(0xFF, 0xF7, 0xED)
C_CREW_LINE = RGBColor(0xF9, 0x73, 0x16)
C_AGENT_FILL = RGBColor(0xFF, 0xED, 0xD5)
C_AGENT_LINE = RGBColor(0xEA, 0x58, 0x0C)
C_AUDIT_FILL = RGBColor(0xE7, 0xFB, 0xE7)
C_AUDIT_LINE = RGBColor(0x16, 0xA3, 0x4A)
C_HITL_FILL = RGBColor(0xF3, 0xF4, 0xF6)
C_HITL_LINE = RGBColor(0x6B, 0x72, 0x80)
C_ARROW = RGBColor(0x4B, 0x55, 0x63)
C_ARROW_FAN = RGBColor(0x9C, 0xA3, 0xAF)  # thinner / muted for fan-out


# ─── 12 specialist agents (one per SGW topic) ─────────────────────────────
SPECIALIST_LABELS = [
    "Financial\nagent",
    "Family\nagent",
    "Caregiving\nagent",
    "Healthcare\nagent",
    "Mental Health\nagent",
    "Crisis\nagent",
    "Disability\nagent",
    "Youth\nagent",
    "Education\nagent",
    "Housing\nagent",
    "Senior\nagent",
    "Employment\nagent",
]


# ─── Helpers ───────────────────────────────────────────────────────────────

def _set_text(shape, text, *, size=11, bold=False, color=C_TEXT,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
              italic=False):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    lines = text.split("\n")
    run = p.add_run(); run.text = lines[0]
    run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color
    for line in lines[1:]:
        ep = tf.add_paragraph(); ep.alignment = align
        r2 = ep.add_run(); r2.text = line
        r2.font.size = Pt(size); r2.font.bold = bold; r2.font.italic = italic
        r2.font.color.rgb = color


def _box(slide, x, y, w, h, fill, line, text, *, size=10, bold=False,
         shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    sh = slide.shapes.add_shape(shape_type, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line; sh.line.width = Pt(1.25)
    sh.shadow.inherit = False
    _set_text(sh, text, size=size, bold=bold)
    return sh


def _arrow(slide, x1, y1, x2, y2, *, color=C_ARROW, width=1.5,
           head_size="med"):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2),
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    ln = conn.line._get_or_add_ln()
    tailEnd = ln.find(qn("a:tailEnd"))
    if tailEnd is None:
        tailEnd = etree.SubElement(ln, qn("a:tailEnd"))
    tailEnd.set("type", "triangle")
    tailEnd.set("w", head_size); tailEnd.set("len", head_size)
    return conn


def _annotation(slide, x, y, w, h, text, *, color=C_MUTED, size=9,
                italic=True, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.color.rgb = color
    run.font.italic = italic
    return tb


def _remove_last_slide_if_matches(prs):
    """If the last slide contains SLIDE_MARKER, remove it."""
    if not prs.slides:
        return False
    last = prs.slides[-1]
    found = False
    for shape in last.shapes:
        if shape.has_text_frame and SLIDE_MARKER in shape.text_frame.text:
            found = True; break
    if not found:
        return False
    idx = len(prs.slides) - 1
    sldIdLst = prs.slides._sldIdLst
    sldId = list(sldIdLst)[idx]
    rId = sldId.attrib[qn("r:id")]
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)
    print(f"Removed previous CrewAI slide (was slide {idx + 1}).")
    return True


# ─── Slide build ───────────────────────────────────────────────────────────

def build_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    SW = prs.slide_width / 914400  # slide width in inches
    CX = SW / 2

    # ─── Title + subtitle ─────────────────────────────────────────────
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.20), Inches(SW - 1.0), Inches(0.50),
    )
    _set_text(title_box, SLIDE_TITLE, size=22, bold=True, color=C_TITLE,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.72), Inches(SW - 1.0), Inches(0.30),
    )
    _set_text(
        subtitle_box,
        "Opt-in multi-agent crew alongside fast single-shot RAG. "
        "Mirrors MSF multidisciplinary case-conference practice.",
        size=11, color=C_MUTED, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
    )

    # ─── Top stack: SAO note → CLOAK → Safety ────────────────────────
    BOX_W = 3.0
    BOX_X = CX - BOX_W / 2

    _box(slide, BOX_X, 1.10, BOX_W, 0.38, C_HITL_FILL, C_HITL_LINE,
         "SAO interview note", size=10, bold=True)
    _arrow(slide, CX, 1.48, CX, 1.60)

    _box(slide, BOX_X, 1.60, BOX_W, 0.42, C_UNCHANGED_FILL, C_UNCHANGED_LINE,
         "Stage 0 — CLOAK PII Guard\n(unchanged · Topic 5.5.2)",
         size=9, bold=True)
    _arrow(slide, CX, 2.02, CX, 2.14)

    _box(slide, BOX_X, 2.14, BOX_W, 0.38, C_UNCHANGED_FILL, C_UNCHANGED_LINE,
         "Safety check (unchanged · Topic 2.6)", size=10, bold=True)
    _arrow(slide, CX, 2.52, CX, 2.64)

    # ─── CrewAI container ─────────────────────────────────────────────
    CR_X = 0.4
    CR_Y = 2.64
    CR_W = SW - 0.8        # 12.53"
    CR_H = 3.55            # 2.64 to 6.19
    container = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(CR_X), Inches(CR_Y), Inches(CR_W), Inches(CR_H),
    )
    container.fill.solid(); container.fill.fore_color.rgb = C_CREW_FILL
    container.line.color.rgb = C_CREW_LINE; container.line.width = Pt(1.5)
    container.shadow.inherit = False

    header = slide.shapes.add_textbox(
        Inches(CR_X + 0.15), Inches(CR_Y + 0.05),
        Inches(CR_W - 0.3), Inches(0.28),
    )
    _set_text(
        header, "CrewAI Hierarchical Crew (Topic 5.5)",
        size=10, bold=True, color=C_CREW_LINE, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP,
    )

    # Triaging Agent (was COORDINATOR/Manager)
    COORD_W = 3.8
    COORD_X = CX - COORD_W / 2
    COORD_Y = CR_Y + 0.40
    COORD_H = 0.55
    _box(slide, COORD_X, COORD_Y, COORD_W, COORD_H,
         C_AGENT_FILL, C_AGENT_LINE,
         "Triaging Agent\nreads case · picks 2–4 of 12 categories",
         size=10, bold=True)
    _annotation(slide, COORD_X + COORD_W + 0.15, COORD_Y + 0.10,
                3.0, 0.4, "← 1 LLM call (gpt-4.1-mini)")

    # ─── 12 specialist agents in a 6 × 2 grid ─────────────────────────
    NCOLS = 6
    NROWS = 2
    SPEC_W = 1.80
    SPEC_GAP_X = 0.18
    SPEC_GAP_Y = 0.12
    SPEC_H = 0.50
    grid_w = NCOLS * SPEC_W + (NCOLS - 1) * SPEC_GAP_X
    SPEC_START_X = CX - grid_w / 2
    SPEC_START_Y = COORD_Y + COORD_H + 0.25       # 1.20 below coord top

    coord_bot_x = CX
    coord_bot_y = COORD_Y + COORD_H

    # Draw the 12 agent boxes + thin fan-out arrows from Coord
    agent_centres = []
    for idx, label in enumerate(SPECIALIST_LABELS):
        col = idx % NCOLS
        row = idx // NCOLS
        ax = SPEC_START_X + col * (SPEC_W + SPEC_GAP_X)
        ay = SPEC_START_Y + row * (SPEC_H + SPEC_GAP_Y)
        _box(slide, ax, ay, SPEC_W, SPEC_H,
             C_AGENT_FILL, C_AGENT_LINE, label, size=9, bold=True)
        cx = ax + SPEC_W / 2
        cy_top = ay
        cy_bot = ay + SPEC_H
        agent_centres.append((cx, cy_top, cy_bot))
        # Thin arrow from Coordinator down to this agent's top edge
        _arrow(slide, coord_bot_x, coord_bot_y, cx, cy_top,
               color=C_ARROW_FAN, width=0.5, head_size="sm")

    # Tool annotation just below the specialist grid
    grid_bottom_y = SPEC_START_Y + NROWS * SPEC_H + (NROWS - 1) * SPEC_GAP_Y
    _annotation(
        slide, SPEC_START_X, grid_bottom_y + 0.05, grid_w, 0.22,
        "Each agent uses a category-filtered retriever as a CrewAI Tool "
        "(dense + HyDE + BM25 + RRF) — Topic 5.5 §Tools",
        align=PP_ALIGN.CENTER,
    )

    # Aggregator Agent (#14)
    AGG_W = 3.6
    AGG_X = CX - AGG_W - 0.25
    AGG_Y = grid_bottom_y + 0.35
    AGG_H = 0.55
    _box(slide, AGG_X, AGG_Y, AGG_W, AGG_H,
         C_AGENT_FILL, C_AGENT_LINE,
         "AGGREGATOR Agent (#14)\nranks top-5 with verbatim citations",
         size=10, bold=True)

    # Case Documentation Officer Agent (#15) — receives Aggregator's output
    DOC_W = 3.6
    DOC_X = CX + 0.25
    DOC_Y = AGG_Y
    DOC_H = 0.55
    _box(slide, DOC_X, DOC_Y, DOC_W, DOC_H,
         C_AGENT_FILL, C_AGENT_LINE,
         "CASE DOCUMENTATION OFFICER Agent (#15)\n"
         "plain-English summary (family + record)",
         size=9, bold=True)

    # Right-side annotation for both bottom-row agents
    _annotation(slide, DOC_X + DOC_W + 0.15, AGG_Y + 0.10,
                2.5, 0.55, "← 2 LLM calls\n(Aggregator, then Docs)")

    # Thin converge arrows from each specialist agent to Aggregator
    agg_top_x = AGG_X + AGG_W / 2
    agg_top_y = AGG_Y
    for cx, _, cy_bot in agent_centres:
        _arrow(slide, cx, cy_bot, agg_top_x, agg_top_y,
               color=C_ARROW_FAN, width=0.5, head_size="sm")

    # Arrow from Aggregator to Case Documentation Officer
    _arrow(
        slide,
        AGG_X + AGG_W, AGG_Y + AGG_H / 2,
        DOC_X, DOC_Y + DOC_H / 2,
        color=C_ARROW, width=1.2, head_size="med",
    )

    # Right-side note: only triaged subset actually executes
    _annotation(
        slide,
        SPEC_START_X + grid_w + 0.15, SPEC_START_Y + 0.05,
        2.8, 0.85,
        "Only the 2–4 agents\npicked by the Triaging\nAgent actually run\n"
        "(in parallel).",
    )

    # ─── Below crew container: Audit → HITL ──────────────────────────
    container_bottom = CR_Y + CR_H
    AUDIT_Y = container_bottom + 0.18
    _arrow(slide, CX, container_bottom, CX, AUDIT_Y)

    _box(slide, BOX_X, AUDIT_Y, BOX_W, 0.38,
         C_AUDIT_FILL, C_AUDIT_LINE,
         "Faithfulness Audit (Topic 4.4)", size=10, bold=True)
    HITL_Y = AUDIT_Y + 0.50
    _arrow(slide, CX, AUDIT_Y + 0.38, CX, HITL_Y)

    _box(slide, BOX_X, HITL_Y, BOX_W, 0.38,
         C_HITL_FILL, C_HITL_LINE,
         "SAO reviews (HITL — end-of-crew)", size=10, bold=True)

    # Footer
    foot = slide.shapes.add_textbox(
        Inches(0.5), Inches(7.10), Inches(SW - 1.0), Inches(0.30),
    )
    _set_text(
        foot,
        "Fast single-shot RAG (v1.3) remains the production default. "
        "Deep Mode opt-in via sidebar toggle. ~$0.02/query · ~25-30 sec · "
        "Topic 5.5 mastery demonstrated.",
        size=9, color=C_MUTED, align=PP_ALIGN.CENTER, italic=True,
    )


def main():
    prs = Presentation(PPT_FILE)
    before = len(prs.slides)
    print(f"Opened PPT — {before} slides before.")
    _remove_last_slide_if_matches(prs)
    build_slide(prs)
    prs.save(PPT_FILE)
    print(f"Saved — {len(prs.slides)} slides after.")


if __name__ == "__main__":
    main()
