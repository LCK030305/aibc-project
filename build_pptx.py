"""Generate the SAO Co-Pilot file-map PowerPoint deck.

Produces SAO_CoPilot_File_Map.pptx in the project root.

Bootcamp-spirit note: this is a build script that runs once (or whenever
the file map changes); keeping it as code (not a manual deck) means the
mapping stays in sync with the project as it evolves.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Palette — "Midnight Executive" with topic-coded category accents.
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x1E, 0x27, 0x61)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x1A, 0x1F, 0x36)
MUTED = RGBColor(0x6B, 0x72, 0x80)
LIGHT_ROW = RGBColor(0xF7, 0xF8, 0xFC)
BORDER = RGBColor(0xD8, 0xDD, 0xE9)

# Category badge colours (last column of every stage table).
CAT_CORE = RGBColor(0xE7, 0x6F, 0x51)      # warm coral — the heart of the runtime
CAT_PIPE = RGBColor(0x2A, 0x9D, 0x8F)      # teal — pipeline / build-time
CAT_DIAG = RGBColor(0x8D, 0x99, 0xAE)      # slate — diagnostics
CAT_CONF = RGBColor(0x26, 0x46, 0x53)      # deep teal-green — config & data

CATEGORY_COLORS: dict[str, RGBColor] = {
    "Core application": CAT_CORE,
    "Data-pipeline scripts": CAT_PIPE,
    "Diagnostics / one-off probes": CAT_DIAG,
    "Config & data": CAT_CONF,
}

# ---------------------------------------------------------------------------
# Presentation set-up.
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK_LAYOUT = prs.slide_layouts[6]


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def add_blank_slide():
    return prs.slides.add_slide(BLANK_LAYOUT)


def add_background(slide, color: RGBColor) -> None:
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()


def add_text(
    slide,
    text: str,
    x: float, y: float, w: float, h: float,
    *,
    font_size: int = 14,
    bold: bool = False,
    color: RGBColor = TEXT_DARK,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    font_name: str = "Calibri",
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, text: str, *, subtitle: str | None = None) -> None:
    add_text(slide, text, 0.5, 0.35, 12.3, 0.7,
             font_size=28, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, subtitle, 0.5, 1.0, 12.3, 0.5,
                 font_size=13, color=MUTED)


def add_stage_table(
    slide,
    headers: list[str],
    rows: list[list[str]],
    *,
    x: float, y: float, w: float, h: float,
    col_widths_pct: list[float],
):
    nrows = len(rows) + 1
    ncols = len(headers)
    tbl_shape = slide.shapes.add_table(
        nrows, ncols, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    table = tbl_shape.table

    # Column widths.
    total_emu = int(w * 914400)
    for i, pct in enumerate(col_widths_pct):
        table.columns[i].width = Emu(int(total_emu * pct))

    # Row heights — header thinner, data rows uniform.
    header_h = Inches(0.45)
    data_h = Inches((h - 0.45) / len(rows))
    table.rows[0].height = header_h
    for r in range(1, nrows):
        table.rows[r].height = data_h

    # Header row.
    for j, label in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        _set_cell_margins(cell)
        tf = cell.text_frame
        tf.word_wrap = True
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = label
        run.font.name = "Calibri"
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = WHITE

    # Data rows.
    for i, row in enumerate(rows, start=1):
        for j, value in enumerate(row):
            cell = table.cell(i, j)
            _set_cell_margins(cell)
            tf = cell.text_frame
            tf.word_wrap = True
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(value)
            run.font.name = "Calibri"
            run.font.size = Pt(10.5)

            is_category_col = j == ncols - 1 and value in CATEGORY_COLORS
            if is_category_col:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CATEGORY_COLORS[value]
                run.font.color.rgb = WHITE
                run.font.bold = True
            else:
                cell.fill.solid()
                # Alternate row tinting for readability.
                cell.fill.fore_color.rgb = WHITE if i % 2 else LIGHT_ROW
                run.font.color.rgb = TEXT_DARK
    return table


def _set_cell_margins(cell) -> None:
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.05)
    cell.margin_bottom = Inches(0.05)


# ---------------------------------------------------------------------------
# Slide 1 — Title
# ---------------------------------------------------------------------------
s = add_blank_slide()
add_background(s, NAVY)

add_text(s, "SAO Co-Pilot", 1.0, 2.0, 11.3, 1.2,
         font_size=56, bold=True, color=WHITE)
add_text(s, "File Map & Bootcamp Alignment", 1.0, 3.2, 11.3, 0.7,
         font_size=26, color=ICE)
# Accent rule
rule = s.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(4.1), Inches(2.0), Inches(0.06)
)
rule.fill.solid()
rule.fill.fore_color.rgb = CAT_CORE
rule.line.fill.background()
add_text(
    s,
    "A Streamlit + RAG capstone for Singapore's Ministry of Social and Family Development. "
    "Each file is mapped to a stage in the pipeline and to the bootcamp topics it embodies.",
    1.0, 4.3, 11.3, 1.5,
    font_size=14, color=ICE,
)
add_text(s, "MSF · Singapore Polytechnic AI Champions Bootcamp · 2026",
         1.0, 6.7, 11.3, 0.4, font_size=12, color=ICE)


# ---------------------------------------------------------------------------
# Slide 2 — Legend (category cards)
# ---------------------------------------------------------------------------
s = add_blank_slide()
add_title(s, "Legend — File categories",
          subtitle="Each row in the next slides is tagged with one of these four roles.")

CARD_W = 5.95
CARD_H = 2.4
CARD_GAP_X = 0.45
CARD_GAP_Y = 0.35
ROW1_Y = 1.7
ROW2_Y = ROW1_Y + CARD_H + CARD_GAP_Y
COL1_X = 0.5
COL2_X = COL1_X + CARD_W + CARD_GAP_X

cards = [
    (COL1_X, ROW1_Y, "Core application", CAT_CORE,
     "Modules used on every recommendation request: LLM wrapper, prompt "
     "templates, retriever, recommender, Streamlit UI. The runtime path."),
    (COL2_X, ROW1_Y, "Data-pipeline scripts", CAT_PIPE,
     "Run once (or on a schedule) to build the SGW corpus: scrape, chunk, "
     "embed, topic-map, patch missing items. Read-then-write."),
    (COL1_X, ROW2_Y, "Diagnostics / one-off probes", CAT_DIAG,
     "Forensics scripts kept in the repo for transparency: DOM probes, gap "
     "diagnoses, sample inspector. Help auditors trust the pipeline."),
    (COL2_X, ROW2_Y, "Config & data", CAT_CONF,
     "Environment + persisted artefacts: requirements.txt, .env (gitignored), "
     ".gitignore, the data/ tree with raw JSON, chunks, embeddings."),
]

for x, y, title, color, body in cards:
    card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                              Inches(CARD_W), Inches(CARD_H))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = BORDER
    card.line.width = Pt(0.75)
    # Coloured left bar.
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                             Inches(0.18), Inches(CARD_H))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_text(s, title, x + 0.35, y + 0.18, CARD_W - 0.5, 0.55,
             font_size=18, bold=True, color=color)
    add_text(s, body, x + 0.35, y + 0.85, CARD_W - 0.5, CARD_H - 1.0,
             font_size=13, color=TEXT_DARK)


# ---------------------------------------------------------------------------
# Common headers + column widths for every stage table.
# ---------------------------------------------------------------------------
STAGE_HEADERS = ["#", "File / Action", "What it does", "Bootcamp topic(s)", "Category"]
STAGE_COL_WIDTHS = [0.06, 0.20, 0.32, 0.22, 0.20]


# ---------------------------------------------------------------------------
# Slide 3 — Stage 0
# ---------------------------------------------------------------------------
s = add_blank_slide()
add_title(s, "Stage 0 · Project setup",
          subtitle="Environment, secrets and version control — laid down before anything else runs.")
add_stage_table(s, STAGE_HEADERS, [
    ["0a", "requirements.txt",
     "Pip deps pinned: playwright, openai, python-dotenv, numpy, streamlit.",
     "Topic 6.2 — pip + venv", "Config & data"],
    ["0b", ".venv/  (python -m venv)",
     "Isolated Python environment so global Python stays untouched.",
     "Topic 6.2 — pip + venv", "Config & data"],
    ["0c", ".env  (gitignored)",
     "OPENAI_API_KEY=… loaded by python-dotenv at import time.",
     "Topic 5.2 — Secure credentials", "Config & data"],
    ["0d", ".gitignore",
     "Excludes .venv/, .env, __pycache__/ from version control.",
     "Topic 8.3 — Intro to Git & GitHub", "Config & data"],
], x=0.5, y=1.6, w=12.3, h=5.0, col_widths_pct=STAGE_COL_WIDTHS)


# ---------------------------------------------------------------------------
# Slide 4 — Stage 1
# ---------------------------------------------------------------------------
s = add_blank_slide()
add_title(s, "Stage 1 · Data acquisition",
          subtitle="One-time build of the SupportGoWhere corpus: 307 records / 2,147 section chunks / 12 category mappings.")
add_stage_table(s, STAGE_HEADERS, [
    ["1", "scraper.py",
     "Playwright renders 296 /schemes/ + /services/ pages; saves cleaned text per record.",
     "Topic 5.3 — Python scripts · Topic 2.7 — Exception handling",
     "Data-pipeline scripts"],
    ["2", "chunker.py",
     "Splits each record into section-level chunks (tagline / highlights / who / apply / help / expect / whole).",
     "Topic 4.2 — Pre-Retrieval · Topic 1.3 — Delimiters",
     "Data-pipeline scripts"],
    ["3", "embed.py",
     "Embeds all 2,147 chunks via text-embedding-3-small; saves vectors.npy + index.jsonl.",
     "Topic 3.1 — Embeddings · Topic 3.2 — Handling Embeddings",
     "Data-pipeline scripts"],
    ["4", "topic_mapper.py",
     "Crawls the 12 /topics/<slug> pages, captures XHR JSON to derive category metadata per record.",
     "Topic 4.2 — Pre-Retrieval (metadata enrichment)",
     "Data-pipeline scripts"],
    ["4b", "scrape_missing.py",
     "Patch scraper for items present on topic pages but absent from the original sitemap (11 found).",
     "Topic 4.5 — RAG Evaluation (coverage check)",
     "Data-pipeline scripts"],
], x=0.5, y=1.6, w=12.3, h=5.3, col_widths_pct=STAGE_COL_WIDTHS)


# ---------------------------------------------------------------------------
# Slide 5 — Stage 2 (Application layer)
# ---------------------------------------------------------------------------
s = add_blank_slide()
add_title(s, "Stage 2 · Application layer",
          subtitle="Runtime modules called on every recommendation request. Designed to be imported anywhere — CLI, Streamlit, future API.")
add_stage_table(s, STAGE_HEADERS, [
    ["5", "llm.py",
     "Thin OpenAI wrapper: get_completion() (chat) + embed_batch() (vectors); loads key from .env.",
     "Topic 1.3 — f-strings · Topic 5.2 — Credentials · Topic 2.7 — Exceptions",
     "Core application"],
    ["6", "prompts.py",
     "CO-STAR template builder; XML-tag delimited candidate rendering.",
     "Topic 1.2 — Prompt Engineering · Topic 1.3 — Delimiters · Topic 2.5 — Multi-action · Playbook p.26 (CO-STAR)",
     "Core application"],
    ["7", "retriever.py",
     "Loads vectors + index + topic map; search() does cosine + dedup + filter; returns typed Result objects.",
     "Topic 3.3 — Applying Embeddings · Topic 3.4 — RAG · Topic 4.3 — Improving Retrieval",
     "Core application"],
    ["8", "recommender.py",
     "End-to-end UC#1: retrieve → render → generate → parse; returns RecommendationResponse.",
     "Topic 2.6 — Prompt Chaining · Topic 4.4 — Post-Retrieval · Topic 2.7 — Exceptions",
     "Core application"],
], x=0.5, y=1.6, w=12.3, h=5.0, col_widths_pct=STAGE_COL_WIDTHS)


# ---------------------------------------------------------------------------
# Slide 6 — Stage 3 (UI)
# ---------------------------------------------------------------------------
s = add_blank_slide()
add_title(s, "Stage 3 · UI",
          subtitle="The SAO-facing surface. Built lightweight so we can iterate fast and deploy to Streamlit Community Cloud later.")
add_stage_table(s, STAGE_HEADERS, [
    ["9", "app.py",
     "Streamlit UI: sidebar filters (category / kind / pool size), sample queries, recommendation cards with rationale + evidence + eligibility flags, debug panels.",
     "Topic 6.1 — Understanding Streamlit · Topic 6.3 — Working with Streamlit · Topic 8.1 — Streamlit Deep Dive (state, @st.cache_resource)",
     "Core application"],
], x=0.5, y=1.6, w=12.3, h=3.0, col_widths_pct=STAGE_COL_WIDTHS)

# Add an aside note since this slide has only one row — fill the space with
# a "What the user sees" callout.
note_y = 5.0
note = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(note_y),
                          Inches(12.3), Inches(1.9))
note.fill.solid()
note.fill.fore_color.rgb = LIGHT_ROW
note.line.color.rgb = BORDER
note.line.width = Pt(0.5)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(note_y),
                         Inches(0.15), Inches(1.9))
bar.fill.solid()
bar.fill.fore_color.rgb = CAT_CORE
bar.line.fill.background()
add_text(s, "What the SAO sees at runtime", 0.85, note_y + 0.1, 11.5, 0.45,
         font_size=15, bold=True, color=CAT_CORE)
add_text(
    s,
    "Enter a 1–2 sentence client situation  →  optional filters  →  "
    "Top-N recommendations with fit-score, plain-English rationale, "
    "verbatim evidence quote from SGW, and a checklist of items the SAO "
    "should verify with the client. Each card links straight to the live "
    "SupportGoWhere page.",
    0.85, note_y + 0.55, 11.5, 1.3,
    font_size=13, color=TEXT_DARK,
)


# ---------------------------------------------------------------------------
# Slide 7 — Stage 4 (Diagnostics)
# ---------------------------------------------------------------------------
s = add_blank_slide()
add_title(s, "Stage 4 · Diagnostics & probes",
          subtitle="One-off scripts kept in the repo so any reviewer can re-run the forensics that informed our design decisions.")
add_stage_table(s, STAGE_HEADERS, [
    ["D1", "probe_dom.py",
     "Initial DOM-structure investigation (homepage + scheme + service pages).",
     "Topic 8.5 — Vibe coding (exploratory probing)",
     "Diagnostics / one-off probes"],
    ["D2", "probe_topic.py",
     "Topic-page DOM probe (revealed that cards are JS-routed buttons, not anchors).",
     "Topic 8.5 — Vibe coding",
     "Diagnostics / one-off probes"],
    ["D3", "diagnose_topic.py · find_missing.py",
     "Per-topic gap diagnosis — surfaced 11 items present on SGW but absent from our sitemap.",
     "Topic 4.5 — RAG Evaluation (coverage)",
     "Diagnostics / one-off probes"],
    ["D4", "inspect_corpus.py",
     "Pretty-print sample records for eyeball QA on chunk + scrape quality.",
     "Topic 4.5 — RAG Evaluation",
     "Diagnostics / one-off probes"],
], x=0.5, y=1.6, w=12.3, h=5.0, col_widths_pct=STAGE_COL_WIDTHS)


# ---------------------------------------------------------------------------
# Slide 8 — Stage 5 (Planned)
# ---------------------------------------------------------------------------
s = add_blank_slide()
add_title(s, "Stage 5 · Planned",
          subtitle="Where the project goes next — already aligned to Weeks 4–9 of the bootcamp.")
add_stage_table(s, STAGE_HEADERS, [
    ["10", "evaluator.py",
     "Ground-truth pairs (~10 client scenarios → expected schemes) + retrieval@k metric + LLM-judge ranking.",
     "Topic 4.5 — RAG Evaluation",
     "Core application"],
    ["11", "doc_parser.py · rules_extractor.py · eligibility_checker.py",
     "UC#2 (Pain Point #2) — extract structured facts from uploaded client docs, match against the 'who' / 'apply' chunks already in our corpus.",
     "Topic 2.5 — Multi-action · Topic 2.6 — Chaining · Topic 4.4 — Post-Retrieval",
     "Core application"],
    ["12", "Streamlit Cloud deployment",
     "Push repo to GitHub, set OPENAI_API_KEY as a Streamlit secret, optional password protection.",
     "Topic 8.2 — Password Protect · Topic 8.3 — Git / GitHub · Topic 8.4 — Deploy to Community Cloud",
     "Config & data"],
], x=0.5, y=1.6, w=12.3, h=5.0, col_widths_pct=STAGE_COL_WIDTHS)


# ---------------------------------------------------------------------------
# Slide 9 — Rubric alignment summary
# ---------------------------------------------------------------------------
s = add_blank_slide()
add_title(s, "Bootcamp curriculum coverage",
          subtitle="Every topic block the bootcamp teaches has a corresponding artefact in the repo.")

ALIGN_ROWS = [
    ["Week 1", "LLM foundations · Prompt Engineering · f-strings · Tokens · Hallucinations",
     "llm.py · prompts.py (CO-STAR)"],
    ["Week 2", "Advanced prompting · Chaining · Multi-action · Exception handling",
     "recommender.py (chained pipeline + JSON parse guard) · llm.py"],
    ["Week 3 (Topic 3.x)", "Embeddings · Handling embeddings · RAG · Search beyond keywords",
     "embed.py · retriever.py"],
    ["Week 4 (Topic 4.x)", "Advanced RAG · Pre / Post-retrieval · Evaluation",
     "chunker.py (pre) · recommender.py (post) · planned evaluator.py"],
    ["Week 5 (Topic 5.x)", "Agents · Secure credentials · Python scripts",
     ".env / python-dotenv · scraper.py · clean script structure"],
    ["Week 6 (Topic 6.x)", "Streamlit basics · pip + venv · Debugging",
     "requirements.txt · .venv/ · app.py"],
    ["Week 8 (Topic 8.x)", "Streamlit Deep Dive · Password protect · Git · Deploy · (Vibe coding)",
     "app.py (@st.cache_resource, session_state) · planned deployment"],
]
add_stage_table(
    s,
    ["Week / Topic", "What the bootcamp teaches", "Where it lives in this repo"],
    ALIGN_ROWS,
    x=0.5, y=1.6, w=12.3, h=5.2,
    col_widths_pct=[0.16, 0.45, 0.39],
)

# Footer
add_text(s, "Capstone deadline: 14 August 2026 · Politemall", 0.5, 7.0, 12.3, 0.4,
         font_size=11, color=MUTED, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# Save
# ---------------------------------------------------------------------------
OUT_PATH = Path(__file__).parent / "SAO_CoPilot_File_Map.pptx"
prs.save(OUT_PATH)
print(f"Wrote: {OUT_PATH}")
print(f"Slides: {len(prs.slides)}")
